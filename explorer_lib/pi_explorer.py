import re


import olefile                  # 한글
from pptx import Presentation   # ppt 
from docx import Document       # 워드
import openpyxl as oxl          # 엑셀
import pdfplumber               # pdf


class PI_Explorer():
    def __init__(self):
        pass
        

    def find_pi(self, file_name, file_type, regexp_li):
        if 'Hangul' in file_type: 
            result = self.__find_pi_hwp(file_name)
            
            return result

        elif 'doc' in file_type:
            result = self.__find_pi_msword(file_name)
            
            
            return result

        elif 'pdf' in file_type:
            
            result = self.__find_pi_pdf(file_name)
            
            
            return result

        elif 'excel' in file_type or 'csv' in file_type:
            result = self.__find_pi_excel(file_name)
                        
            return result

        elif 'ppt' in file_type:
            result = self.__find_pi_ppt(file_name)
            return result

        elif 'txt' in file_type:
            result = self.__find_pi_txt(file_name)
            return result

        else:
            print(f'{file_type}은 지원되지 않는 문서 형식 입니다.')
            return False
        

    def __read_hwp(self, file_name):
        pass


    def __read_word(self, file_name):
        pass


    def __read_pdf(self, file_name):
        pass


    def __read_excel(self, file_name):
        pass

    #보류
    def __read_ppt(self, file_name):
        pass


    def __read_txt(self, file_name):
        pass

    def __find_pi_hwp(self, file_name):
        """ 한글 문서 내 문자열 추출

        Args:
            file_name (str): 한글 문서의 경로

        Returns:
            list: 문서 내 텍스트 데이터

        TODO:
            * 표 내의 글 읽기
            * 이미지 읽기 -> OCR 처리
            
        """
        file = olefile.OleFileIO(file_name)
        
        encoded_data = file.openstream('PrvText').read()
        decoded_data = encoded_data.decode('utf-16')

        # img = file.openstream('PrvImage').read()
        return decoded_data 


    def __find_pi_msword(self, file_name):
        """ doc, docx 파일 내 문자열 추출
        
        Args:
            file_name (str): 문서의 경로

        Returns:
            list: 문서 내의 모든 텍스트

        TODO:
            * 이미지 읽기 -> OCR처리
        """
        data = []
        doc = Document(file_name)
        text_data = [paragraph.text for paragraph in doc.paragraphs]

        tables = doc.tables
        tables_data = []
        table_data = [cell.text for table in tables for row in table.rows for cell in row.cells]
        data.append(text_data for txt_data in text_data)
        data.append(tb_data for tb_data in table_data)

        return data
        

    def __find_pi_pdf(self, file_name):
        
        pdf = pdfplumber.open(file_name)
        pages = pdf.pages
        data =[]

        for page in pages:
            data.append((lambda page: page.extract_text())(page))

        return data

    def __find_pi_excel(self, file_name):
        """_summary_

        Args:
            file_name (_type_): _description_

        Returns:
            _type_: _description_
            
        TODO:
            * 엑셀 시트 순회 구현 

        """
        wb = oxl.load_workbook(file_name, data_only=True)
        ws = wb.active

        all_values = [[cell.value for cell in row]for row in ws.rows]   
        data = sum(all_values, [])                                # 1차원 리스트로 변환
        
        # print(all_values)
        return data

    
    def __find_pi_ppt(self, file_name):
        data = []
        prs = Presentation(file_name)


        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    data.append(paragraph.text)
        
        return data


    def __find_pi_txt(self, file_name):
        pass

    def find_pi(file_obj, regexp):
        pass

if __name__ == '__main__':
    file_name = './test/test.pptx'
    file_type ='pptx'
    regexp_li = ''

    

    


    px = PI_Explorer()
    
    temp = px.find_pi(file_name, file_type, regexp_li)    
